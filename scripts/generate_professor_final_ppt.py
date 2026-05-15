from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("outputs/presentations/professor_final_deck.pptx")


BG = RGBColor(255, 255, 255)
TEXT = RGBColor(17, 17, 17)
MUTED = RGBColor(88, 88, 88)
ACCENT = RGBColor(38, 84, 124)
ACCENT_LIGHT = RGBColor(232, 240, 247)
BORDER = RGBColor(206, 214, 224)
GOOD = RGBColor(58, 117, 75)
WARN = RGBColor(155, 85, 0)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(0.7))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT
    return box


def add_bullets(slide, items, left, top, width, height, level0=20, level1=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    frame.clear()
    first = True
    for item in items:
        if first:
            p = frame.paragraphs[0]
            first = False
        else:
            p = frame.add_paragraph()
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p.text = text
        p.level = level
        p.font.name = "Helvetica Neue"
        p.font.size = Pt(level0 if level == 0 else level1)
        p.font.color.rgb = TEXT if level == 0 else MUTED
        p.space_after = Pt(6)
        p.bullet = True
    return box


def add_text_block(slide, text, left, top, width, height, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_mono_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT
    shape.line.color.rgb = BORDER
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Menlo"
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT
        p.space_after = Pt(0)
    return shape


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12.0), Inches(0.35))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(8.5)
    run.font.color.rgb = MUTED
    return box


def add_header_rule(slide):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(1.05), Inches(12.6), Inches(1.05))
    line.line.color.rgb = BORDER
    line.line.width = Pt(1.25)


def add_labeled_box(slide, text, left, top, width, height, fill_color=BG, line_color=BORDER, font_size=15, bold=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = TEXT
    return shape


def connect(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.75)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Self-Evolving Image Editing")
    add_header_rule(slide)

    add_text_block(slide, "Known pieces already exist", Inches(0.7), Inches(1.25), Inches(5.3), Inches(0.35), size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Self-evolving proposer-solver training exists for multimodal reasoning.",
            "RL and reward-based post-training already exist for image editing.",
            "Specialized edit scorers and reward models already exist.",
        ],
        Inches(0.7), Inches(1.6), Inches(5.7), Inches(2.0)
    )

    add_text_block(slide, "Actual gap", Inches(6.75), Inches(1.25), Inches(5.0), Inches(0.35), size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Image editing needs both requested change and preservation of everything else.",
            "Existing editing RL usually optimizes one editor against one reward model.",
            "The missing combination is a proposer-editor-solver curriculum over raw unlabeled images.",
        ],
        Inches(6.75), Inches(1.6), Inches(5.4), Inches(2.0)
    )

    claim = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.25), Inches(12.0), Inches(1.35))
    claim.fill.solid()
    claim.fill.fore_color.rgb = ACCENT_LIGHT
    claim.line.color.rgb = BORDER
    tf = claim.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Claim: a self-evolving image editing loop where the proposer generates edit instructions on raw images, "
        "the editor samples multiple candidates, and the solver filters them using editing-specific "
        "constraint-based and relative rewards."
    )
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(19)
    run.font.bold = True
    run.font.color.rgb = TEXT

    add_text_block(
        slide,
        'Not generic RL or generic self-play; the novelty is an editing-specific verifier that separates requested change from collateral damage.',
        Inches(0.7), Inches(5.95), Inches(12.0), Inches(0.65), size=18, color=WARN, bold=False
    )

    add_footer(
        slide,
        "Refs: EvoLMM (2511.16672) | SQLM (2508.03682) | MM-Zero (2603.09206) | InstructRL4Pix (2406.09973) | EditReward (2509.26346) | ADIEE (2507.07317)"
    )


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Final Architecture")
    add_header_rule(slide)

    add_mono_box(
        slide,
        "Proposer -> Editor(K=4) -> Solver Ensemble -> Relative Ranker -> Accept/Reject -> Train",
        Inches(0.7), Inches(1.2), Inches(12.0), Inches(0.7)
    )

    proposer = add_labeled_box(slide, "Proposer\nGenerate edit instruction", Inches(0.8), Inches(2.2), Inches(1.9), Inches(1.0), fill_color=ACCENT_LIGHT, line_color=ACCENT, font_size=15, bold=True)
    raw = add_labeled_box(slide, "Raw image x", Inches(0.8), Inches(3.55), Inches(1.9), Inches(0.75), font_size=16)
    editor = add_labeled_box(slide, "Editor\nQwen-Image-Edit\nSample K=4 candidates", Inches(3.15), Inches(2.65), Inches(2.35), Inches(1.35), fill_color=ACCENT_LIGHT, line_color=ACCENT, font_size=15, bold=True)
    c1 = add_labeled_box(slide, "y1", Inches(6.05), Inches(1.95), Inches(0.9), Inches(0.6), font_size=16)
    c2 = add_labeled_box(slide, "y2", Inches(6.05), Inches(2.75), Inches(0.9), Inches(0.6), font_size=16)
    c3 = add_labeled_box(slide, "y3", Inches(6.05), Inches(3.55), Inches(0.9), Inches(0.6), font_size=16)
    c4 = add_labeled_box(slide, "y4", Inches(6.05), Inches(4.35), Inches(0.9), Inches(0.6), font_size=16)
    solver = add_labeled_box(slide, "Solver ensemble\ninst | pres | spatial | semantic", Inches(7.55), Inches(2.55), Inches(2.35), Inches(1.55), fill_color=ACCENT_LIGHT, line_color=ACCENT, font_size=14, bold=True)
    ranker = add_labeled_box(slide, "Relative ranker", Inches(10.35), Inches(2.8), Inches(1.65), Inches(0.75), fill_color=ACCENT_LIGHT, line_color=ACCENT, font_size=15, bold=True)
    accept = add_labeled_box(slide, "Accept / Reject", Inches(10.25), Inches(4.1), Inches(1.85), Inches(0.75), fill_color=ACCENT_LIGHT, line_color=ACCENT, font_size=15, bold=True)
    train = add_labeled_box(slide, "Pseudo-label pool\nTrain next round", Inches(8.75), Inches(5.45), Inches(3.2), Inches(0.95), fill_color=ACCENT_LIGHT, line_color=GOOD, font_size=15, bold=True)

    connect(slide, Inches(1.75), Inches(3.2), Inches(1.75), Inches(3.55))
    connect(slide, Inches(2.7), Inches(2.7), Inches(3.15), Inches(2.7))
    connect(slide, Inches(2.7), Inches(3.9), Inches(3.15), Inches(3.5))
    connect(slide, Inches(5.5), Inches(3.3), Inches(6.05), Inches(2.25))
    connect(slide, Inches(5.5), Inches(3.3), Inches(6.05), Inches(3.05))
    connect(slide, Inches(5.5), Inches(3.3), Inches(6.05), Inches(3.85))
    connect(slide, Inches(5.5), Inches(3.3), Inches(6.05), Inches(4.65))
    for y in [Inches(2.25), Inches(3.05), Inches(3.85), Inches(4.65)]:
        connect(slide, Inches(6.95), y, Inches(7.55), Inches(3.3))
    connect(slide, Inches(9.9), Inches(3.3), Inches(10.35), Inches(3.15))
    connect(slide, Inches(11.15), Inches(3.55), Inches(11.15), Inches(4.1))
    connect(slide, Inches(11.1), Inches(4.85), Inches(10.45), Inches(5.45))

    add_bullets(
        slide,
        [
            "K-sample editing matters because image editing is multimodal.",
            "The ranker compares candidates instead of trusting one absolute score.",
            "This is the EvoLMM logic adapted to a visual transformation task.",
        ],
        Inches(0.8), Inches(5.45), Inches(7.1), Inches(1.2), level0=15
    )

    add_footer(slide, "Refs: EvoLMM (2511.16672) | SQLM (2508.03682) | MM-Zero (2603.09206)")


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Final Reward Functions")
    add_header_rule(slide)

    add_text_block(slide, "1. Hard-gated feasibility", Inches(0.7), Inches(1.2), Inches(4.0), Inches(0.35), size=18, color=ACCENT, bold=True)
    add_mono_box(slide, "G(y) = 1[S_inst(x,e,y) >= tau_inst] * 1[S_pres(x,y) >= tau_pres]", Inches(0.7), Inches(1.55), Inches(5.65), Inches(0.85))
    add_bullets(
        slide,
        [
            "S_inst: did the requested edit happen?",
            "S_pres: did unchanged content stay unchanged?",
            "These are constraints, not soft preferences.",
        ],
        Inches(0.7), Inches(2.55), Inches(5.6), Inches(1.4), level0=15
    )

    add_text_block(slide, "2. Relative quality score", Inches(6.7), Inches(1.2), Inches(4.0), Inches(0.35), size=18, color=ACCENT, bold=True)
    add_mono_box(
        slide,
        "Q(y) = alpha * S_spa(x,e,y)\n     + beta  * S_cf(x,e+,y,{e-})\n     + gamma * S_rel(x,e,y,y_ref)",
        Inches(6.7), Inches(1.55), Inches(5.55), Inches(1.2)
    )
    add_bullets(
        slide,
        [
            "S_spa: spatial correctness",
            "S_cf: true instruction beats distractors",
            "S_rel: improvement over a reference output",
        ],
        Inches(6.7), Inches(2.95), Inches(5.2), Inches(1.25), level0=15
    )

    add_text_block(slide, "3. Acceptance and proposer reward", Inches(0.7), Inches(4.15), Inches(4.7), Inches(0.35), size=18, color=ACCENT, bold=True)
    add_mono_box(
        slide,
        "accept best y* only if:\nG(y*)=1\nQ(y*) >= tau_q\nstd(component_scores(y*)) <= delta_conf",
        Inches(0.7), Inches(4.5), Inches(5.1), Inches(1.35)
    )
    add_mono_box(
        slide,
        "R_prop = exp(- (u - mu)^2 / (2 * sigma^2))\n       - lambda_fail * 1[all rejected]\n       - lambda_easy * 1[too easy]",
        Inches(6.7), Inches(4.5), Inches(5.45), Inches(1.35)
    )
    add_text_block(
        slide,
        "u = normalized difficulty from candidate disagreement. The solver reward is hard-gated and relative; the proposer reward is band-pass over difficulty.",
        Inches(0.7), Inches(6.15), Inches(11.5), Inches(0.5), size=16, color=WARN
    )

    add_footer(slide, "Refs: EvoLMM (2511.16672) | InstructRL4Pix (2406.09973) | SpatialReward (2602.07458) | EditReward (2509.26346) | MDPO (2406.11839)")


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Why This Could Work")
    add_header_rule(slide)

    add_text_block(slide, "Main intuition", Inches(0.7), Inches(1.2), Inches(3.5), Inches(0.35), size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "A single scalar edit score is too easy to game.",
            "Image editing has two required objectives: requested change and preservation.",
            "Hard constraints stop catastrophic compensation between reward terms.",
            "Relative ranking is more stable than one absolute reward.",
            "The proposer follows the EvoLMM learning-frontier logic.",
        ],
        Inches(0.7), Inches(1.55), Inches(5.8), Inches(2.8), level0=15
    )

    add_text_block(slide, "Main risks", Inches(6.8), Inches(1.2), Inches(3.0), Inches(0.35), size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Reward hacking",
            "Proposer collapse",
            "Solver miscalibration",
        ],
        Inches(6.8), Inches(1.55), Inches(4.6), Inches(1.3), level0=15
    )

    add_text_block(slide, "Experiment ladder", Inches(6.8), Inches(3.1), Inches(3.0), Inches(0.35), size=18, color=ACCENT, bold=True)
    ladder = [
        "1. supervised baseline",
        "2. naive self-training",
        "3. current hybrid reward",
        "4. hard-gated + K-sample ranking",
        "5. add counterfactual reward",
        "6. add reference-relative reward",
    ]
    box = slide.shapes.add_textbox(Inches(6.8), Inches(3.45), Inches(5.0), Inches(2.0))
    tf = box.text_frame
    tf.clear()
    for i, text in enumerate(ladder):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Helvetica Neue"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.space_after = Pt(4)

    status = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.75), Inches(5.8), Inches(1.25))
    status.fill.solid()
    status.fill.fore_color.rgb = ACCENT_LIGHT
    status.line.color.rgb = BORDER
    tf = status.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Status: baseline training/eval stack is implemented; self-evolving loop is implemented; next step is the final K-sample and relative-ranking reward path on the GPU machine."
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(16)
    run.font.color.rgb = TEXT

    add_text_block(
        slide,
        "Key question: does better reward design lead to better self-generated edit data, not just more self-training?",
        Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.45), size=17, color=GOOD, bold=True
    )

    add_footer(slide, "Refs: EvoLMM (2511.16672) | EditReward (2509.26346) | ADIEE (2507.07317) | SpatialReward (2602.07458) | MDPO (2406.11839)")


def main():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
