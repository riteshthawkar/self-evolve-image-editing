from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("outputs/presentations/professor_progress_deck.pptx")

BG = RGBColor(255, 255, 255)
TEXT = RGBColor(17, 17, 17)
MUTED = RGBColor(95, 95, 95)
ACCENT = RGBColor(38, 84, 124)
ACCENT_LIGHT = RGBColor(232, 240, 247)
BORDER = RGBColor(206, 214, 224)
GOOD = RGBColor(58, 117, 75)
WARN = RGBColor(155, 85, 0)
SOFT_WARN = RGBColor(255, 247, 230)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT


def add_header_rule(slide):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.7), Inches(1.05), Inches(12.6), Inches(1.05))
    line.line.color.rgb = BORDER
    line.line.width = Pt(1.25)


def add_text(slide, text, left, top, width, height, size=18, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, left, top, width, height, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Helvetica Neue"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)
        p.bullet = True
    return box


def add_mono_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT
    shape.line.color.rgb = BORDER
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Menlo"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT
        p.space_after = Pt(0)
    return shape


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.92), Inches(12.0), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(8.5)
    run.font.color.rgb = MUTED


def add_note_box(slide, text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SOFT_WARN
    shape.line.color.rgb = WARN
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(14)
    run.font.color.rgb = WARN
    run.font.bold = True
    return shape


def add_table(slide, rows, cols, left, top, width, height, header_fill=ACCENT_LIGHT):
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text_frame.clear()
    return table


def write_cell(cell, text, bold=False, color=TEXT, size=13, align=PP_ALIGN.LEFT):
    cell.text_frame.clear()
    cell.text_frame.word_wrap = True
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)


def color_cell(cell, fill_rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_rgb


def add_architecture_strip(slide, left, top):
    def box(text, x, w, fill=ACCENT_LIGHT):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, w, Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = BORDER
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Helvetica Neue"
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = TEXT
        return shape

    x = left
    b1 = box("Proposer", x, Inches(1.6)); x += Inches(1.9)
    b2 = box("Editor (K=4)", x, Inches(1.9)); x += Inches(2.2)
    b3 = box("Solver Ensemble", x, Inches(2.25)); x += Inches(2.55)
    b4 = box("Relative Ranker", x, Inches(2.0)); x += Inches(2.3)
    b5 = box("Accept / Train", x, Inches(1.9), fill=RGBColor(233, 244, 235))
    boxes = [b1, b2, b3, b4, b5]
    for i in range(len(boxes)-1):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            boxes[i].left + boxes[i].width,
            boxes[i].top + boxes[i].height/2,
            boxes[i+1].left,
            boxes[i+1].top + boxes[i+1].height/2,
        )
        line.line.color.rgb = ACCENT
        line.line.width = Pt(1.8)
        try:
            line.line.end_arrowhead = True
        except Exception:
            pass


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Current Status and Final Direction")
    add_header_rule(slide)
    add_architecture_strip(slide, Inches(0.85), Inches(1.35))

    add_bullets(
        slide,
        [
            "The broad problem is fixed: self-evolving image editing from raw images.",
            "The method is narrowed to one final architecture and reward design.",
            "The project is no longer at the brainstorming stage; infrastructure and local validation are already done.",
        ],
        Inches(0.8), Inches(2.55), Inches(5.8), Inches(2.0)
    )

    add_text(slide, "What changed", Inches(6.95), Inches(2.55), Inches(2.5), Inches(0.3), size=18, bold=True, color=ACCENT)
    add_bullets(
        slide,
        [
            "Started from a simple weighted hybrid reward.",
            "After checking EvoLMM-style proposer curriculum and recent edit reward papers, the final method became:",
            "hard-gated instruction and preservation checks, relative ranking over K candidates, optional counterfactual and reference-relative rewards.",
        ],
        Inches(6.95), Inches(2.9), Inches(5.1), Inches(2.2)
    )

    add_note_box(
        slide,
        "Meeting framing: show completed engineering work and method narrowing first; show benchmark numbers later only when measured.",
        Inches(0.8), Inches(5.55), Inches(11.2), Inches(0.9)
    )

    add_footer(slide, "Refs: EvoLMM (2511.16672) | SQLM (2508.03682) | MM-Zero (2603.09206) | InstructRL4Pix (2406.09973) | EditReward (2509.26346)")


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Work Already Completed")
    add_header_rule(slide)

    table = add_table(slide, 8, 3, Inches(0.7), Inches(1.45), Inches(12.0), Inches(4.75))
    headers = ["Component", "Status", "Evidence"]
    for i, h in enumerate(headers):
        write_cell(table.cell(0, i), h, bold=True, size=14)

    rows = [
        ("Baseline training stack", "Implemented", "LoRA and full training launchers"),
        ("Editing evaluation", "Implemented", "GEdit and ImgEdit export plus scoring"),
        ("Generation sanity evaluation", "Implemented", "GenEval, DPG-Bench, OneIG runners"),
        ("Self-evolving loop", "Implemented", "proposer-editor-solver loop and configs"),
        ("Ablation variants", "Implemented", "spatial, cycle, internal, hybrid configs"),
        ("Resume and run tooling", "Implemented", "resumable shell runners and smoke tests"),
        ("Local pipeline validation", "Completed", "compile checks, dry runs, pillow demo"),
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            write_cell(table.cell(r, c), text, size=13)
        color_cell(table.cell(r, 1), RGBColor(236, 245, 237))

    add_text(
        slide,
        "Message: the remaining step is running the full GPU experiments, not building the project from scratch.",
        Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.4), size=16, color=GOOD, bold=True
    )

    add_footer(slide, "Actual completed work shown above; no benchmark claims on this slide.")


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Methods Already Checked")
    add_header_rule(slide)

    add_note_box(
        slide,
        "These values are conservative planning estimates, not measured benchmark results.",
        Inches(0.8), Inches(1.2), Inches(11.0), Inches(0.6)
    )

    table = add_table(slide, 8, 5, Inches(0.7), Inches(1.95), Inches(12.0), Inches(4.15))
    headers = ["Method", "Why I checked it", "Est. GEdit-EN", "Est. ImgEdit", "Current read"]
    for i, h in enumerate(headers):
        write_cell(table.cell(0, i), h, bold=True, size=14)

    rows = [
        ("Qwen-Image official baseline", "public baseline reference", "7.56", "4.27", "comparison anchor"),
        ("Plain proxy reward", "simplest self-training baseline", "7.22", "4.08", "too easy to game conceptually"),
        ("Weighted hybrid reward", "combine edit success with preservation signals", "7.86", "4.46", "better, but still allows bad compensation"),
        ("Spatial-only ablation", "test localization signal alone", "7.74", "4.33", "useful, but incomplete alone"),
        ("Cycle-only ablation", "test reversibility signal alone", "7.66", "4.24", "stabilizing, but likely too restrictive alone"),
        ("Internal-only ablation", "test semantic internal verifier alone", "7.70", "4.29", "promising, but probably too weak alone"),
        ("Hard-gated + relative ranker", "stronger anti-hacking and better fit to multimodal editing", "8.08", "4.66", "best current paper direction"),
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            write_cell(table.cell(r, c), text, size=12.5)
    color_cell(table.cell(1, 4), RGBColor(233, 240, 250))
    color_cell(table.cell(7, 4), RGBColor(233, 244, 235))

    add_text(
        slide,
        "This slide is about method narrowing, not benchmark victory. It shows that multiple reward designs were already considered and filtered.",
        Inches(0.8), Inches(6.28), Inches(11.0), Inches(0.35), size=14, color=WARN, bold=True
    )
    add_footer(slide, "Refs: EvoLMM (2511.16672) | InstructRL4Pix (2406.09973) | SpatialReward (2602.07458) | EditReward (2509.26346) | MDPO (2406.11839)")


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Experiment Matrix and Placeholder Result Table")
    add_header_rule(slide)

    add_note_box(
        slide,
        "Important: values below are planning targets or placeholders, not measured results. Replace them after GPU runs.",
        Inches(0.8), Inches(1.25), Inches(11.4), Inches(0.75)
    )

    table = add_table(slide, 7, 5, Inches(0.7), Inches(2.0), Inches(12.0), Inches(3.9))
    headers = ["Variant", "Run status", "GEdit-EN", "ImgEdit", "Note"]
    for i, h in enumerate(headers):
        write_cell(table.cell(0, i), h, bold=True, size=13)

    rows = [
        ("Qwen-Image official baseline", "official", "7.56", "4.27", "public report value"),
        ("Our local supervised baseline", "estimate", "7.30", "4.12", "conservative local reproduction estimate"),
        ("Naive self-training", "estimate", "7.18", "4.06", "mainly a sanity check"),
        ("Current hybrid reward", "target", "7.8 - 8.2", "4.4 - 4.7", "expected modest gain if filtering helps"),
        ("Final hard-gated + ranker", "next main run", "8.0 - 8.6", "4.6 - 5.0", "current main method"),
        ("+ counterfactual reward", "optional follow-up", "8.2 - 8.8", "4.7 - 5.1", "optional extension"),
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            write_cell(table.cell(r, c), text, size=12.5)
        if r in [5, 6]:
            color_cell(table.cell(r, 1), RGBColor(233, 244, 235))
        elif r in [1]:
            color_cell(table.cell(r, 1), RGBColor(233, 240, 250))
        elif r in [4]:
            color_cell(table.cell(r, 1), RGBColor(255, 249, 235))
        elif r in [2, 3]:
            color_cell(table.cell(r, 1), RGBColor(244, 239, 255))

    add_bullets(
        slide,
        [
            "Do not present any estimate or target range as a result.",
            "Use this slide to show a clear experiment ladder and that the search space is already reduced.",
            "After the first GPU runs, replace all non-official rows with measured benchmark numbers.",
        ],
        Inches(0.8), Inches(6.15), Inches(11.4), Inches(0.8), size=15
    )
    add_footer(slide, "Use this slide only if you clearly say: placeholder targets, actual results pending.")


def main():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
